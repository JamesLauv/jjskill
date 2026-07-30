"""
PPTX内容提取脚本
从PowerPoint文件中提取幻灯片文本、表格和备注。

用法: python extract_pptx.py <file.pptx> [--json]
输出: JSON格式的幻灯片内容（默认）或纯文本
"""
import sys
import json

def extract_pptx(filepath):
    from pptx import Presentation
    prs = Presentation(filepath)
    slides_data = []

    for i, slide in enumerate(prs.slides, 1):
        slide_info = {
            "slide_number": i,
            "texts": [],
            "tables": [],
            "notes": ""
        }

        for shape in slide.shapes:
            # 提取文本框内容
            if hasattr(shape, "text") and shape.text.strip():
                slide_info["texts"].append(shape.text.strip())

            # 提取表格内容
            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    table_data.append(cells)
                slide_info["tables"].append(table_data)

        # 提取备注（在slide级别检查，不在shape级别）
        try:
            if slide.has_notes_slide and slide.notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_info["notes"] = notes_text
        except (AttributeError, TypeError):
            pass

        slides_data.append(slide_info)

    return slides_data


def print_text(slides_data):
    """纯文本输出格式"""
    import io
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

    for slide in slides_data:
        print(f"\n{'='*60}")
        print(f"=== Slide {slide['slide_number']} ===")
        print(f"{'='*60}")
        for text in slide["texts"]:
            print(text)
        for ti, table in enumerate(slide["tables"], 1):
            print(f"\n--- Table {ti} ---")
            for row in table:
                print(" | ".join(row))
        if slide["notes"]:
            print(f"\n[Notes] {slide['notes']}")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("用法: python extract_pptx.py <file.pptx> [--json]")
        sys.exit(1)

    filepath = sys.argv[1]
    use_json = "--json" in sys.argv

    try:
        data = extract_pptx(filepath)
        if use_json:
            import io
            sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
            print(json.dumps(data, ensure_ascii=False, indent=2))
        else:
            print_text(data)
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)
